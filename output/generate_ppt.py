from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

SCREENSHOTS = os.path.join(os.path.dirname(__file__), "screenshots")
OUT = os.path.join(os.path.dirname(__file__), "MultiModal_GraphRAG.pptx")

BG    = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACC   = RGBColor(0xAA, 0xAA, 0xAA)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text_box(slide, txt, left, top, width, height,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def slide_title(slide, title, subtitle=None):
    text_box(slide, title,
             Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
             size=36, bold=True, color=WHITE)
    if subtitle:
        text_box(slide, subtitle,
                 Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                 size=16, color=GRAY)
    # divider line
    from pptx.util import Pt as UPt
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.fill.background()


def add_image(slide, filename, left, top, width):
    path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width=width)


def bullet_box(slide, items, left, top, width, height, size=17):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = LGRAY


def table_slide(slide, headers, rows, left, top, width, col_widths=None):
    from pptx.util import Inches
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, left, top, width, Inches(0.4 * nrows)).table
    tbl.first_row = True

    def cell_text(cell, txt, bold=False, color=WHITE, size=14):
        cell.text = txt
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x11) if not bold else RGBColor(0x1a, 0x1a, 0x1a)

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    for ci, h in enumerate(headers):
        cell_text(tbl.cell(0, ci), h, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell_text(tbl.cell(ri + 1, ci), val, color=LGRAY)


# ── SLIDE 1 — TITLE ─────────────────────────────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "MultiModal Graph RAG",
         Inches(0.6), Inches(1.2), Inches(12), Inches(1.2),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text_box(s, "A Document Intelligence Platform — Chat · Upload · Visualize",
         Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
text_box(s, '"Ask questions across all your files — PDFs, images, audio, videos\n— and get intelligent, sourced answers in real time."',
         Inches(1.5), Inches(3.2), Inches(10), Inches(1.0),
         size=16, color=ACC, align=PP_ALIGN.CENTER)
text_box(s, "FastAPI  ·  ChromaDB  ·  Neo4j  ·  Groq AI  ·  React  ·  Docker",
         Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
# Team info
text_box(s, "Team: DATA MESH",
         Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text_box(s, "Thanmai A  ·  1MS23CI127          Varun F Kshatriya  ·  1MS24CI411",
         Inches(0.6), Inches(5.45), Inches(12), Inches(0.4),
         size=14, color=LGRAY, align=PP_ALIGN.CENTER)
add_image(s, "chat_interface_premium.png", Inches(2.7), Inches(5.95), Inches(7.8))


# ── SLIDE 2 — THE PROBLEM ────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "The Problem", "Why existing tools fall short")
table_slide(s,
    ["File Type", "Traditional Tools", "This System"],
    [
        ["PDFs & documents",      "Can read",   "Can read"],
        ["Images & photos",       "Ignored",    "AI vision — captioned + searchable"],
        ["Audio files",           "Ignored",    "Whisper transcription + search"],
        ["Videos",                "Ignored",    "Audio extracted, transcribed, indexed"],
        ["Cross-file connections","No memory",  "Knowledge graph connects everything"],
    ],
    Inches(0.6), Inches(1.5), Inches(12.1),
    col_widths=[Inches(3), Inches(3.2), Inches(5.9)]
)
text_box(s,
    "Old tools answer from one file at a time.\nThis system answers from ALL your files — connected together.",
    Inches(0.6), Inches(5.6), Inches(12), Inches(0.8),
    size=18, bold=True, color=WHITE)


# ── SLIDE 3 — FEATURES ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "Feature Overview", "What the system can do")

col_left  = Inches(0.6)
col_right = Inches(6.9)
cw        = Inches(6.0)

text_box(s, "File Intelligence", col_left, Inches(1.5), cw, Inches(0.35), size=15, bold=True, color=WHITE)
bullet_box(s, [
    "•  PDF, TXT, MD, JPG, PNG, MP3, WAV, MP4, MOV — all formats",
    "•  AI vision captions every image (Groq LLaMA 4 Scout)",
    "•  Whisper transcribes every audio and video file",
    "•  SHA-256 deduplication — same file never indexed twice",
], col_left, Inches(1.85), cw, Inches(1.1), size=14)

text_box(s, "Smart Chat", col_left, Inches(3.05), cw, Inches(0.35), size=15, bold=True, color=WHITE)
bullet_box(s, [
    "•  Plain English queries — typos handled automatically",
    "•  Streaming responses, token by token",
    "•  Clickable source badges — open images, play audio, view PDFs",
    "•  Context-aware follow-ups across conversation turns",
], col_left, Inches(3.4), cw, Inches(1.1), size=14)

text_box(s, "Knowledge Graph + Engineering", col_left, Inches(4.6), cw, Inches(0.35), size=15, bold=True, color=WHITE)
bullet_box(s, [
    "•  Extracts entities and relationships from every file",
    "•  Auto-heals graph when files are deleted",
    "•  Dockerized — 4 services, one command to start",
], col_left, Inches(4.95), cw, Inches(0.9), size=14)

add_image(s, "filemanager_premium.png", col_right, Inches(1.5), Inches(5.8))


# ── SLIDE 4 — ARCHITECTURE ──────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "System Architecture", "How all components connect")
arch = (
    "Browser (React)\n"
    "  Chat  ·  Upload  ·  Knowledge Graph\n"
    "         │   HTTP + SSE streaming\n"
    "         ▼\n"
    "FastAPI Backend  (Python 3.11)\n"
    "  Upload → Process → Index → Search\n"
    "    │                       │\n"
    "    ▼                       ▼\n"
    "ChromaDB                 Neo4j\n"
    "Vector Store          Knowledge Graph\n"
    "Stores meaning        Stores connections\n"
    "as embeddings         between concepts"
)
text_box(s, arch, Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.5),
         size=14, color=LGRAY)

text_box(s, "ChromaDB", Inches(7.4), Inches(1.9), Inches(5.3), Inches(0.4), size=16, bold=True, color=WHITE)
text_box(s, "Finds relevant content.\n'What chunks talk about this topic?'",
         Inches(7.4), Inches(2.35), Inches(5.3), Inches(0.8), size=14, color=LGRAY)

text_box(s, "Neo4j", Inches(7.4), Inches(3.4), Inches(5.3), Inches(0.4), size=16, bold=True, color=WHITE)
text_box(s, "Finds connected knowledge.\n'What else in the system relates to this?'",
         Inches(7.4), Inches(3.85), Inches(5.3), Inches(0.8), size=14, color=LGRAY)

text_box(s, "Together they produce richer, more accurate answers than either alone.",
         Inches(7.4), Inches(5.1), Inches(5.3), Inches(0.7), size=14, bold=True, color=WHITE)


# ── SLIDE 5 — PIPELINE ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "File Processing Pipeline", "From upload to searchable in seconds")
pipeline = (
    "User uploads a file\n"
    "        │\n"
    "        ▼\n"
    "Safety Check\n"
    "  · Reject duplicates (SHA-256)\n"
    "  · Reject files over 50 MB\n"
    "        │\n"
    "   ┌────┴────────────────┐\n"
    "   ▼                     ▼\n"
    "TEXT / PDF            IMAGE\n"
    "· Extract text        · CLIP embedding\n"
    "· Chunk + embed       · Vision AI caption\n"
    "· Store ChromaDB      · Store image collection\n"
    "                  AUDIO / VIDEO\n"
    "                  · FFmpeg extracts audio\n"
    "                  · Whisper transcribes\n"
    "                  · Treated as text above\n"
    "        │\n"
    "        ▼\n"
    "Knowledge Graph Extraction\n"
    "  · LLaMA 3.3 extracts entities + relationships\n"
    "  · Stored in Neo4j as nodes and edges"
)
text_box(s, pipeline, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.8), size=13, color=LGRAY)
text_box(s, "Result",
         Inches(8.2), Inches(1.9), Inches(4.7), Inches(0.4), size=16, bold=True, color=WHITE)
text_box(s,
    "Every file — regardless of type — becomes searchable text and graph knowledge within seconds.",
    Inches(8.2), Inches(2.35), Inches(4.7), Inches(1.2), size=14, color=LGRAY)
add_image(s, "filemanager_premium.png", Inches(8.2), Inches(3.8), Inches(4.7))


# ── SLIDE 6 — CHAT SYSTEM ───────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "The Chat System", "How the AI answers a question end to end")
flow = (
    'User types: "Explain the attention mechanism"\n'
    "        │\n"
    "   ┌────┴─────────────────────┐\n"
    "   ▼                          ▼\n"
    "Text embedding (MiniLM)    CLIP embedding (images)\n"
    "   │                          │\n"
    "   ▼                          ▼\n"
    "Search ChromaDB text       Search ChromaDB images\n"
    "→ Relevant PDF chunks      → Relevant images\n"
    "   │                          │\n"
    "   └──────────────┬───────────┘\n"
    "                  │\n"
    "                  ▼\n"
    "       Extract keywords → Query Neo4j\n"
    "       Pull related graph entities\n"
    "                  │\n"
    "                  ▼\n"
    "       Combine into context (12k chars max)\n"
    "                  │\n"
    "                  ▼\n"
    "       LLaMA 3.3 70B on Groq\n"
    "       Stream answer token by token"
)
text_box(s, flow, Inches(0.6), Inches(1.5), Inches(7.0), Inches(5.8), size=12, color=LGRAY)
add_image(s, "chat_interface_premium.png", Inches(7.7), Inches(1.5), Inches(5.3))
text_box(s, "Answer streams with [Source 1] [Source 2] badges — click any to open the file.",
         Inches(7.7), Inches(6.2), Inches(5.3), Inches(0.6), size=13, color=GRAY)


# ── SLIDE 7 — SMART RETRIEVAL ────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "Smart Retrieval", "Handling real-world messy input")

text_box(s, "1. Typo-Tolerant Modality Detection",
         Inches(0.6), Inches(1.5), Inches(12), Inches(0.4), size=16, bold=True, color=WHITE)
table_slide(s,
    ["What the user types", "System understands", "Result"],
    [
        ['"show me images"',   "image modality",         "Images only"],
        ['"show me vdios"',    "video (fuzzy match)",    "Videos only"],
        ['"giv me aduio"',     "audio (fuzzy match)",    "Audio only"],
        ['"what pdfs do I have"', "pdf format filter",   "PDFs only"],
    ],
    Inches(0.6), Inches(1.95), Inches(12.1),
    col_widths=[Inches(3.5), Inches(3.8), Inches(4.8)]
)
text_box(s, "2. Context-Aware Follow-ups",
         Inches(0.6), Inches(3.7), Inches(12), Inches(0.4), size=16, bold=True, color=WHITE)
table_slide(s,
    ["Previous question", "Follow-up", "What happens"],
    [
        ['"show me city.jpg"',         '"tell me about it"', "Retrieves city.jpg"],
        ['"give me sample_audio.mp3"', '"play it"',          "Returns that audio file"],
    ],
    Inches(0.6), Inches(4.15), Inches(12.1),
    col_widths=[Inches(4.0), Inches(3.5), Inches(4.6)]
)
text_box(s, "3. Filename-Specific Search",
         Inches(0.6), Inches(5.5), Inches(8), Inches(0.4), size=16, bold=True, color=WHITE)
text_box(s, 'Mention any filename — full or partial — and the system returns exactly that file. "city" matches "city.jpg".',
         Inches(0.6), Inches(5.9), Inches(12), Inches(0.5), size=14, color=LGRAY)


# ── SLIDE 8 — KNOWLEDGE GRAPH ────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "Knowledge Graph", "Connecting knowledge across all files")
example = (
    "Example — upload bitcoin_whitepaper.pdf + ai_overview.txt:\n\n"
    "  [Satoshi Nakamoto] ── CREATED ──→ [Bitcoin]\n"
    "          [Bitcoin] ── IS_A ──────→ [Cryptocurrency]\n"
    "       [Blockchain] ── ENABLES ──→ [Bitcoin]\n"
    "      [Transformer] ── USED_IN ──→ [LLaMA 3.3]\n"
    "  [Attention Paper] ── INTRODUCED→ [Transformer]\n"
)
text_box(s, example, Inches(0.6), Inches(1.5), Inches(6.5), Inches(2.5), size=13, color=LGRAY)

text_box(s, "Auto-Heal on Delete", Inches(0.6), Inches(4.0), Inches(6.5), Inches(0.4), size=16, bold=True, color=WHITE)
bullet_box(s, [
    "1.  File's vector chunks and entities are removed",
    "2.  Disconnected neighbors are identified",
    "3.  LLaMA evaluates whether a direct relationship makes sense",
    "4.  If yes — link is created automatically, graph stays intact",
], Inches(0.6), Inches(4.45), Inches(6.5), Inches(1.5), size=14)

add_image(s, "graph_premium.png", Inches(7.1), Inches(1.5), Inches(5.9))


# ── SLIDE 9 — TECH STACK ────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "Tech Stack", "All free and open source")
table_slide(s,
    ["Layer", "Technology", "Purpose"],
    [
        ["Frontend",     "React 18 + Vite + Tailwind CSS",   "Chat, file manager, graph UI"],
        ["Backend",      "FastAPI (Python 3.11)",             "API, processing, orchestration"],
        ["Vector DB",    "ChromaDB",                         "Semantic similarity search"],
        ["Graph DB",     "Neo4j 5.15",                       "Entity relationship storage"],
        ["LLM",          "Groq — LLaMA 3.3 70B",             "Chat responses, entity extraction"],
        ["Vision AI",    "Groq — LLaMA 4 Scout 17B",         "Image captioning"],
        ["Speech AI",    "Groq — Whisper Large v3",           "Audio and video transcription"],
        ["Embeddings",   "sentence-transformers MiniLM",      "Text to vector conversion"],
        ["Image Search", "CLIP ViT-B/32",                    "Cross-modal image search"],
        ["Containers",   "Docker + Docker Compose",           "One-command deployment"],
    ],
    Inches(0.6), Inches(1.5), Inches(12.1),
    col_widths=[Inches(1.9), Inches(4.0), Inches(6.2)]
)
text_box(s, "cp .env.example .env  &&  ./start.sh  →  http://localhost:3000",
         Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 10 — RESULTS ──────────────────────────────────────────────────────
s = add_slide(); bg(s)
slide_title(s, "Results & Takeaways", "What was built and what was learned")

col_l = Inches(0.6)
col_r = Inches(7.0)
cw2   = Inches(5.9)

text_box(s, "Working Features", col_l, Inches(1.5), cw2, Inches(0.4), size=16, bold=True, color=WHITE)
bullet_box(s, [
    "•  4 file modalities — text, image, audio, video",
    "•  Streaming chat with source citations, under 2 seconds",
    "•  Typo-tolerant queries via fuzzy matching",
    "•  Context memory — follow-ups resolve naturally",
    "•  Knowledge graph with auto-heal on file deletion",
    "•  Parallel uploads — multiple files indexed at once",
    "•  One-command Docker deployment on any machine",
], col_l, Inches(1.95), cw2, Inches(2.5), size=14)

text_box(s, "Scale", col_l, Inches(4.6), cw2, Inches(0.4), size=16, bold=True, color=WHITE)
bullet_box(s, [
    "•  Hundreds of files across all modalities",
    "•  Thousands of vector chunks in ChromaDB",
    "•  Thousands of entity nodes in Neo4j",
], col_l, Inches(5.0), cw2, Inches(1.0), size=14)

text_box(s, "Key Takeaway", col_r, Inches(1.5), cw2, Inches(0.4), size=16, bold=True, color=WHITE)
text_box(s,
    "This is not just a chatbot.\n\n"
    "It is a complete document intelligence system — built from scratch — "
    "that understands every file type, connects knowledge across documents, "
    "and runs entirely on free APIs.",
    col_r, Inches(1.95), cw2, Inches(2.5), size=16, color=WHITE)

add_image(s, "chat_interface_premium.png", col_r, Inches(4.6), cw2)


# ── SAVE ────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
